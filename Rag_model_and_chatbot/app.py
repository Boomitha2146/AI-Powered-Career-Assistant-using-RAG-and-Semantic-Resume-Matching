
import os
import streamlit as st
import pandas as pd
from datasets import load_dataset
from sentence_transformers import SentenceTransformer, util
import docx2txt
from pptx import Presentation
import fitz  # PyMuPDF

from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline

# LangChain (LATEST SAFE IMPORTS)
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.llms.huggingface_pipeline import HuggingFacePipeline
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_groq import ChatGroq

# ---------------- CONFIG ----------------
st.set_page_config(
    page_title="AI Career Assistant",
    page_icon="💼",
    layout="wide"
)

BASE_PATH = os.getcwd()
FAISS_PATH = os.path.join(BASE_PATH, "faiss_index")

# ---------------- CACHE MODELS ----------------
@st.cache_resource(show_spinner=False)
def load_resume_model():
    return SentenceTransformer("all-MiniLM-L6-v2", device="cpu")


@st.cache_resource(show_spinner=False)
def load_embeddings_and_faiss():
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/paraphrase-MiniLM-L3-v2"
    )

    if not os.path.exists(FAISS_PATH):
        return embeddings, None

    vectorstore = FAISS.load_local(
        FAISS_PATH,
        embeddings=embeddings,
        allow_dangerous_deserialization=True
    )
    return embeddings, vectorstore


@st.cache_resource(show_spinner=False)
def load_flan_llm():
    model_name = "google/flan-t5-small"
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)

    pipe = pipeline(
        "text2text-generation",
        model=model,
        tokenizer=tokenizer,
        max_new_tokens=200,
        device=-1
    )
    return HuggingFacePipeline(pipeline=pipe)


@st.cache_resource(show_spinner=False)
def build_rag_chain():
    _, vectorstore = load_embeddings_and_faiss()

    if vectorstore is None:
        return None

    llm = ChatGroq(
        groq_api_key=st.secrets["GROQ_API_KEY"],
        model_name="llama-3.1-8b-instant",
        temperature=0.2
    )

    prompt = PromptTemplate.from_template("""
You are an expert IT career assistant.
Answer using ONLY the given context.
If the answer is not present, say "I don't have that information."

Context:
{context}

Question:
{question}

Answer:
""")

    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})

    rag_chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    return rag_chain


# ---------------- DATA ----------------
@st.cache_data
def load_data():
    dataset = load_dataset("DevilsLord/It_job_roles_skills_certifications")
    df = dataset["train"].to_pandas()

    df["combined"] = (
        df["Job Description"].fillna("") + " " +
        df["Skills"].fillna("") + " " +
        df["Certifications"].fillna("")
    )
    return df


# ---------------- UTILS ----------------
def extract_text(file):
    ext = file.name.split(".")[-1].lower()

    if ext == "txt":
        return file.read().decode("utf-8", errors="ignore")

    if ext == "pdf":
        doc = fitz.open(stream=file.read(), filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    if ext == "docx":
        return docx2txt.process(file)

    if ext == "pptx":
        text = ""
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return ""


# ---------------- INIT ----------------
with st.spinner("Loading AI models..."):
    df = load_data()
    resume_model = load_resume_model()
    rag_chain = build_rag_chain()


# ---------------- UI ----------------
st.sidebar.title("Navigation")
page = st.sidebar.radio("", ["🏠 Home", "📄 Resume Matcher", "💬 Chatbot"])

# ---------------- HOME ----------------
if page == "🏠 Home":
    st.title("🧠 AI-Powered Career Assistant")
    st.write("Match your resume & ask career questions using AI.")


# ---------------- RESUME MATCHER ----------------
elif page == "📄 Resume Matcher":
    st.title("📄 Resume Matcher")

    uploaded_file = st.file_uploader(
        "Upload Resume (PDF / DOCX / PPTX / TXT)",
        type=["pdf", "docx", "pptx", "txt"]
    )

    if uploaded_file:
        resume_text = extract_text(uploaded_file)

        if resume_text.strip():
            resume_emb = resume_model.encode(resume_text, convert_to_tensor=True)
            job_embs = resume_model.encode(df["combined"].tolist(), convert_to_tensor=True)

            scores = util.cos_sim(resume_emb, job_embs)[0].cpu().numpy()
            df["score"] = scores

            top = df.sort_values("score", ascending=False).head(10)

            for _, row in top.iterrows():
                with st.expander(f"{row['Job Title']} — {row['score']:.2f}"):
                    st.write("**Description:**", row["Job Description"])
                    st.write("**Skills:**", row["Skills"])
                    st.write("**Certifications:**", row["Certifications"])


# ---------------- CHATBOT ----------------
elif page == "💬 Chatbot":
    st.title("💬 Career Chatbot")

    if rag_chain is None:
        st.error("FAISS index not found. Build it first.")
        st.stop()

    if "messages" not in st.session_state:
        st.session_state.messages = [
            {"role": "assistant", "content": "Hi! Ask me about IT careers."}
        ]

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

    if user_input := st.chat_input("Ask a career question..."):
        st.session_state.messages.append({"role": "user", "content": user_input})

        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                answer = rag_chain.invoke(user_input)
                st.write(answer)

        st.session_state.messages.append({"role": "assistant", "content": answer})

           


def main():
  
    pass 

if __name__ == "__main__":
    main()
